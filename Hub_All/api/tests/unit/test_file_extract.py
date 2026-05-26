"""Unit test file_extract — Plan 04-02 Task 01 (INGEST-02 prerequisite).

Test pure Python — KHÔNG cần Postgres/Redis. Sample DOCX/TXT/PDF được
generate runtime trong tmp_path để tránh commit binary fixture vào git.

Coverage 7 case theo plan behavior matrix:
- DOCX VN có heading "Mục N." → text + is_scanned=False + meta.format=docx.
- TXT UTF-8 VN có dấu → text intact + encoding=utf-8.
- TXT non-UTF-8 (latin-1) → chardet fallback + decode.
- PDF text-only → is_scanned=False.
- PDF scanned (page text < 30 char) → is_scanned=True (R4 mitigation).
- Extension không thuộc whitelist → UnsupportedFormatError.
- File không tồn tại → FileNotFoundError.
"""
from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document as DocxDocument
from pypdf import PdfWriter

from app.services.file_extract import (
    ALLOWED_EXTENSIONS,
    UnsupportedFormatError,
    detect_scanned_pdf,
    extract_text,
)

# ---------- Fixtures ----------


@pytest.fixture
def docx_vn(tmp_path: Path) -> Path:
    """Tạo DOCX VN có 2 heading + 2 paragraph nội dung."""
    doc = DocxDocument()
    doc.add_paragraph("Mục 1. KHÁM TỔNG QUÁT")
    doc.add_paragraph("Bệnh nhân được khám lâm sàng đầy đủ.")
    doc.add_paragraph("Mục 2. XÉT NGHIỆM")
    doc.add_paragraph("Làm xét nghiệm máu và nước tiểu.")
    path = tmp_path / "khám-bệnh.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def docx_tables_only(tmp_path: Path) -> Path:
    """DOCX layout-bảng — 0 paragraph nội dung, toàn bộ chữ nằm trong table cell.

    Tái hiện đúng dạng tài liệu nội bộ Medinet (vd PhanCong_NhanVat) dựng layout
    bằng bảng — case mà `_extract_docx` cũ (chỉ đọc paragraphs) extract ra rỗng.
    """
    doc = DocxDocument()
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "PHÂN CÔNG NHÂN VẬT"
    table.cell(0, 1).text = "Mã T3-02"
    table.cell(1, 0).text = "Nội dung quan trọng nằm trong ô bảng."
    table.cell(1, 1).text = "Tư vấn viên phụ trách kênh."
    path = tmp_path / "bảng-layout.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def txt_utf8_vn(tmp_path: Path) -> Path:
    path = tmp_path / "khám-utf8.txt"
    path.write_text("Khám bệnh đa khoa\nLâm sàng tốt.", encoding="utf-8")
    return path


@pytest.fixture
def txt_latin1_vn(tmp_path: Path) -> Path:
    """TXT encode latin-1 — UTF-8 decode sẽ fail → chardet fallback."""
    path = tmp_path / "latin1.txt"
    # latin-1 byte sequence không phải UTF-8 valid (bytes 0xc0-0xff không follow UTF-8).
    path.write_bytes(b"H\xe0 N\xf4i\nL\xe1m s\xe0ng")  # "Hà Nội Lâm sàng" latin-1
    return path


@pytest.fixture
def pdf_text_only(tmp_path: Path) -> Path:
    """Tạo PDF có text-only — pypdf extract_text trả nội dung > 30 char/page.

    Dùng pypdf.PdfWriter.add_blank_page rồi annotate text — nhưng pypdf không
    có API thuận tiện tạo text PDF. Workaround: tạo PDF từ reportlab nếu có,
    hoặc tạo PDF tối giản từ raw bytes.

    Phương án đơn giản nhất: tạo PDF bằng pypdf merge từ existing template.
    Với plan này, để tránh thêm dependency reportlab, ta tạo PDF "minimal valid"
    với 1 page chứa 1 stream text qua dictionary direct.
    """
    # PDF tối thiểu hợp lệ (RFC PDF 1.4) với 1 page text "Khám bệnh đa khoa lâm sàng tốt"
    # Hand-crafted để đảm bảo pypdf parse được + extract_text trả > 30 char.
    pdf_bytes = (
        b"%PDF-1.4\n"
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n"
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n"
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj\n"
        b"4 0 obj << /Length 90 >> stream\n"
        b"BT\n/F1 12 Tf\n50 750 Td\n(Kham benh da khoa lam sang tot va day du noi dung text only) Tj\nET\n"
        b"endstream\nendobj\n"
        b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n"
        b"xref\n0 6\n0000000000 65535 f\n0000000009 00000 n\n0000000058 00000 n\n"
        b"0000000115 00000 n\n0000000232 00000 n\n0000000370 00000 n\n"
        b"trailer << /Size 6 /Root 1 0 R >>\n"
        b"startxref\n442\n%%EOF\n"
    )
    path = tmp_path / "text-only.pdf"
    path.write_bytes(pdf_bytes)
    return path


@pytest.fixture
def pdf_scanned(tmp_path: Path) -> Path:
    """Tạo PDF scanned — pypdf extract_text trả empty/garbage cho mọi page.

    Dùng pypdf.PdfWriter add_blank_page → page không có text content stream
    → extract_text() trả "" → is_scanned=True (heuristic > 80% page < 30 char).
    """
    writer = PdfWriter()
    for _ in range(5):
        writer.add_blank_page(width=612, height=792)
    path = tmp_path / "scanned.pdf"
    with path.open("wb") as f:
        writer.write(f)
    return path


# ---------- Tests ----------


def test_extract_docx_vietnamese(docx_vn: Path) -> None:
    """DOCX VN extract intact với heading + sentence."""
    text, is_scanned, meta = extract_text(docx_vn)
    assert "Mục 1. KHÁM TỔNG QUÁT" in text
    assert "Mục 2. XÉT NGHIỆM" in text
    assert "lâm sàng" in text
    assert is_scanned is False
    assert meta["format"] == "docx"
    assert meta["pages"] == 1
    assert meta["paragraph_count"] == 4


def test_extract_docx_tables(docx_tables_only: Path) -> None:
    """DOCX toàn bảng (0 paragraph) — extract phải ra nội dung trong table cell.

    Regression: `_extract_docx` cũ chỉ đọc `doc.paragraphs` → tài liệu layout-bảng
    extract ra rỗng → 0 chunk → document kẹt. Core value M2 yêu cầu tái hiện "bảng".
    """
    text, is_scanned, meta = extract_text(docx_tables_only)
    assert "PHÂN CÔNG NHÂN VẬT" in text
    assert "Mã T3-02" in text
    assert "Nội dung quan trọng nằm trong ô bảng." in text
    assert "Tư vấn viên phụ trách kênh." in text
    assert is_scanned is False
    assert meta["format"] == "docx"
    assert meta["table_count"] == 1
    assert meta["paragraph_count"] == 0


def test_extract_txt_utf8_vietnamese(txt_utf8_vn: Path) -> None:
    """TXT UTF-8 VN có dấu — intact."""
    text, is_scanned, meta = extract_text(txt_utf8_vn)
    assert "Khám bệnh" in text
    assert "Lâm sàng" in text
    assert is_scanned is False
    assert meta["format"] == "txt"
    assert meta["encoding"] == "utf-8"


def test_extract_txt_latin1_fallback(txt_latin1_vn: Path) -> None:
    """TXT non-UTF-8 → chardet fallback decode (latin-1 / Windows-1258)."""
    text, is_scanned, meta = extract_text(txt_latin1_vn)
    assert text  # KHÔNG empty
    assert is_scanned is False
    assert meta["format"] == "txt"
    # Encoding KHÔNG phải utf-8 — chardet detect được latin-1 hoặc tương đương.
    assert meta["encoding"] != "utf-8"


def test_extract_pdf_text_only(pdf_text_only: Path) -> None:
    """PDF text-only → is_scanned=False, pages > 0."""
    text, is_scanned, meta = extract_text(pdf_text_only)
    assert is_scanned is False, f"PDF text-only KHÔNG scanned. Got text: {text!r}"
    assert meta["pages"] >= 1
    assert meta["format"] == "pdf"


def test_extract_pdf_scanned_detected(pdf_scanned: Path) -> None:
    """PDF 5 page blank → > 80% page < 30 char → is_scanned=True (R4)."""
    text, is_scanned, meta = extract_text(pdf_scanned)
    assert is_scanned is True, (
        f"PDF blank PHẢI detect scanned. text={text!r} pages={meta['pages']}"
    )
    assert meta["pages"] == 5


def test_detect_scanned_pdf_public_api(pdf_scanned: Path, pdf_text_only: Path) -> None:
    """Public detect_scanned_pdf wrapper — Plan 04-04 BLOCKER #3 prerequisite.

    Router phải import detect_scanned_pdf SYNC để reject 415 sớm trước khi
    queue cocoindex flow.
    """
    assert detect_scanned_pdf(pdf_scanned) is True
    assert detect_scanned_pdf(pdf_text_only) is False


def test_extract_unsupported_extension(tmp_path: Path) -> None:
    """Extension .exe không trong whitelist → UnsupportedFormatError."""
    bad = tmp_path / "malware.exe"
    bad.write_bytes(b"MZ")
    with pytest.raises(UnsupportedFormatError) as exc:
        extract_text(bad)
    assert ".exe" in str(exc.value)
    assert exc.value.ext == ".exe"


def test_extract_file_not_found(tmp_path: Path) -> None:
    """File không tồn tại → FileNotFoundError."""
    with pytest.raises(FileNotFoundError):
        extract_text(tmp_path / "nope.docx")


def test_allowed_extensions_pinned() -> None:
    """Whitelist post-v3.1 quick task 2026-05-26: 8 ext (R4 mitigation extend).

    Original M2 4 format + quick task extend 4 format (csv/xlsx/pptx/html).
    OCR ảnh (.jpg/.png) + .doc legacy vẫn out of scope (defer v4.0).
    """
    assert ALLOWED_EXTENSIONS == frozenset(
        {".docx", ".txt", ".md", ".pdf", ".csv", ".xlsx", ".pptx", ".html"}
    )


def test_extract_md_file(tmp_path: Path) -> None:
    """MD file extract như TXT — share code path."""
    path = tmp_path / "doc.md"
    path.write_text("# Khám bệnh\n\nNội dung markdown.", encoding="utf-8")
    text, is_scanned, meta = extract_text(path)
    assert "Khám bệnh" in text
    assert is_scanned is False
    assert meta["format"] == "md"


# ============================================================================
# Quick task 2026-05-26-add-file-format-readers — 4 format mới
# ============================================================================


@pytest.fixture
def csv_vn(tmp_path: Path) -> Path:
    """CSV 3 cột × 4 row VN có dấu, UTF-8 (KHÔNG BOM)."""
    import csv as stdlib_csv

    path = tmp_path / "danh-sach-thuoc.csv"
    with path.open("w", encoding="utf-8", newline="") as f:
        writer = stdlib_csv.writer(f)
        writer.writerow(["Mã thuốc", "Tên thuốc", "Đơn vị"])
        writer.writerow(["T001", "Paracetamol 500mg", "Viên"])
        writer.writerow(["T002", "Amoxicillin 250mg", "Viên"])
        writer.writerow(["T003", "Vitamin C", "Gói"])
    return path


@pytest.fixture
def csv_excel_bom_semicolon(tmp_path: Path) -> Path:
    """CSV xuất từ Excel Windows VN — UTF-8 BOM + delimiter ';'."""
    path = tmp_path / "excel-export.csv"
    content = "﻿Họ tên;Tuổi;Chẩn đoán\nNguyễn Văn A;45;Cảm cúm\nTrần Thị B;30;Viêm họng\n"
    path.write_bytes(content.encode("utf-8"))
    return path


@pytest.fixture
def xlsx_two_sheets(tmp_path: Path) -> Path:
    """XLSX 2 sheet — Bệnh nhân + Thuốc."""
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Bệnh nhân"
    ws1.append(["ID", "Họ tên", "Tuổi"])
    ws1.append([1, "Nguyễn Văn A", 45])
    ws1.append([2, "Trần Thị B", 30])

    ws2 = wb.create_sheet("Thuốc")
    ws2.append(["Mã", "Tên thuốc"])
    ws2.append(["T001", "Paracetamol 500mg"])
    ws2.append(["T002", "Amoxicillin 250mg"])

    path = tmp_path / "ho-so.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def pptx_three_slides(tmp_path: Path) -> Path:
    """PPTX 3 slide — title + content."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[5]  # Title Only layout (luôn có ở default template)

    for i, (title, body) in enumerate(
        [
            ("Mục 1. KHÁM TỔNG QUÁT", "Bệnh nhân khám lâm sàng đầy đủ."),
            ("Mục 2. XÉT NGHIỆM", "Làm xét nghiệm máu và nước tiểu."),
            ("Mục 3. KẾT LUẬN", "Bệnh nhân ổn định."),
        ],
        start=1,
    ):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        # Add 1 text box body
        from pptx.util import Inches

        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        tb.text_frame.text = body
        del i  # unused

    path = tmp_path / "trinh-bay.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def html_with_script(tmp_path: Path) -> Path:
    """HTML có <script> + <style> để verify decompose."""
    path = tmp_path / "trang.html"
    html = (
        "<!DOCTYPE html><html><head>"
        "<title>Khám bệnh</title>"
        "<style>body { color: red; }</style>"
        "<script>alert('xss should be stripped');</script>"
        "</head><body>"
        "<h1>Khám bệnh đa khoa</h1>"
        "<p>Bệnh nhân lâm sàng tốt.</p>"
        "<ul><li>Xét nghiệm máu</li><li>Xét nghiệm nước tiểu</li></ul>"
        "</body></html>"
    )
    path.write_text(html, encoding="utf-8")
    return path


def test_extract_csv_vietnamese(csv_vn: Path) -> None:
    """CSV UTF-8 VN — row join ' | ', detect delimiter ','."""
    text, is_scanned, meta = extract_text(csv_vn)
    assert "Mã thuốc | Tên thuốc | Đơn vị" in text
    assert "Paracetamol 500mg" in text
    assert "Vitamin C" in text
    assert is_scanned is False
    assert meta["format"] == "csv"
    assert meta["row_count"] == 4  # 1 header + 3 data
    assert meta["delimiter"] == ","


def test_extract_csv_excel_bom_semicolon(csv_excel_bom_semicolon: Path) -> None:
    """CSV BOM + delimiter ';' — utf-8-sig strip BOM + Sniffer detect ';'."""
    text, is_scanned, meta = extract_text(csv_excel_bom_semicolon)
    # BOM ﻿ PHẢI bị strip — KHÔNG có ở đầu text decoded
    assert not text.startswith("﻿")
    assert "Nguyễn Văn A" in text
    assert "Cảm cúm" in text
    # Delimiter ';' được Sniffer detect
    assert meta["delimiter"] == ";"
    assert meta["encoding"] == "utf-8-sig"
    assert is_scanned is False


def test_extract_xlsx_two_sheets(xlsx_two_sheets: Path) -> None:
    """XLSX 2 sheet — extract cả sheet với separator '--- Sheet: ... ---'."""
    text, is_scanned, meta = extract_text(xlsx_two_sheets)
    assert "--- Sheet: Bệnh nhân ---" in text
    assert "--- Sheet: Thuốc ---" in text
    assert "Nguyễn Văn A" in text
    assert "Paracetamol 500mg" in text
    assert is_scanned is False
    assert meta["format"] == "xlsx"
    assert meta["sheet_count"] == 2
    # 3 row × 2 sheet = 6 row data (header + 2 data row mỗi sheet)
    assert meta["row_count"] == 6


def test_extract_pptx_three_slides(pptx_three_slides: Path) -> None:
    """PPTX 3 slide — extract title + body với separator '--- Slide N ---'."""
    text, is_scanned, meta = extract_text(pptx_three_slides)
    assert "--- Slide 1 ---" in text
    assert "--- Slide 2 ---" in text
    assert "--- Slide 3 ---" in text
    assert "Mục 1. KHÁM TỔNG QUÁT" in text
    assert "xét nghiệm máu" in text.lower()
    assert "Bệnh nhân ổn định" in text
    assert is_scanned is False
    assert meta["format"] == "pptx"
    assert meta["slide_count"] == 3
    assert meta["pages"] == 3


def test_extract_html_decompose_script_style(html_with_script: Path) -> None:
    """HTML — <script> + <style> decompose + tag stripped + content giữ nguyên."""
    text, is_scanned, meta = extract_text(html_with_script)
    # Content user-visible PHẢI có
    assert "Khám bệnh đa khoa" in text
    assert "Bệnh nhân lâm sàng tốt" in text
    assert "Xét nghiệm máu" in text
    # Script content PHẢI bị strip (T-quick-html-01 mitigation)
    assert "alert" not in text
    assert "xss" not in text
    # Style content PHẢI bị strip
    assert "color: red" not in text
    # Tag stripped — KHÔNG có '<' hoặc '>' raw trong output
    assert "<h1>" not in text
    assert "<p>" not in text
    assert is_scanned is False
    assert meta["format"] == "html"


def test_extract_legacy_doc_still_rejected(tmp_path: Path) -> None:
    """Word 97-2003 .doc binary vẫn out of scope (user phải convert sang .docx)."""
    bad = tmp_path / "old-word.doc"
    bad.write_bytes(b"\xd0\xcf\x11\xe0")  # OLE2 magic header
    with pytest.raises(UnsupportedFormatError) as exc:
        extract_text(bad)
    assert ".doc" in str(exc.value)
    assert exc.value.ext == ".doc"
