"""File extract service — Plan 04-02 Task 01 (INGEST-02 prerequisite).

Hỗ trợ 8 format whitelist (R4 mitigation): DOCX, TXT, MD, PDF, CSV, XLSX, PPTX, HTML.
KHÔNG OCR cho ảnh + scanned PDF (D4 — defer v4.0). Scanned PDF → is_scanned=True.

Public API ổn định cho:
- Plan 04-03 cocoindex flow wrap thành cocoindex.op.function (extract step).
- Plan 04-04 router /api/documents/upload import detect_scanned_pdf để reject 415
  TRƯỚC KHI queue ingest job (BLOCKER #3 strategy A — router-side sync detection).

Quy tắc xử lý theo extension:

| Extension     | Library                | Output                  | is_scanned                  |
|---------------|------------------------|-------------------------|-----------------------------|
| .docx         | python-docx Document() | paragraph + bảng        | False luôn                  |
| .txt, .md     | open + chardet detect  | full text               | False luôn                  |
| .pdf          | pypdf PdfReader        | pages join \\n          | True nếu > 80% page < 30 ký |
| .csv          | stdlib csv + Sniffer   | row join " | "          | False luôn                  |
| .xlsx         | openpyxl read_only     | sheet/row join " | "    | False luôn                  |
| .pptx         | python-pptx            | slide text + heading    | False luôn                  |
| .html         | bs4 + lxml             | get_text strip tag      | False luôn                  |

Tham chiếu:
- PITFALLS P5 — Scanned PDF silent fail (HIGH).
- PROJECT.md R4 — whitelist 4 format M2 (CSV/XLSX/PPTX/HTML thêm post-v3.1 quick task 2026-05-26).
- CLAUDE.md section 3 — format hỗ trợ.
- Quick task 2026-05-26-add-file-format-readers — mở rộng 4 → 8 format.
"""
from __future__ import annotations

import csv as _csv
import io as _io
from collections.abc import Iterator
from pathlib import Path
from typing import Any

import chardet
from docx import Document as DocxDocument
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from pypdf import PdfReader

#: Whitelist extension hỗ trợ — KHÔNG đổi giữa phase (R4 mitigation).
#: Quick task 2026-05-26 mở rộng 4 → 8 (thêm csv/xlsx/pptx/html).
ALLOWED_EXTENSIONS: frozenset[str] = frozenset(
    {".docx", ".txt", ".md", ".pdf", ".csv", ".xlsx", ".pptx", ".html"}
)

#: Threshold heuristic detect scanned PDF (P5 / R4):
#: page có < 30 ký tự non-whitespace coi là "không có text"; > 80% page như vậy
#: → toàn document scanned. 30 ký tự cover edge case "Trang 1" + page number.
_SCANNED_PAGE_CHAR_THRESHOLD: int = 30
_SCANNED_RATIO_THRESHOLD: float = 0.8


class UnsupportedFormatError(Exception):
    """Raise khi extension không nằm trong ALLOWED_EXTENSIONS.

    Plan 04-04 router catch exception này → trả 415 Unsupported Media Type
    với message tiếng Việt khuyến nghị format thay thế.
    """

    def __init__(self, ext: str) -> None:
        super().__init__(
            f"Định dạng {ext!r} không hỗ trợ trong M2. "
            f"Khuyến nghị: chuyển sang một trong {sorted(ALLOWED_EXTENSIONS)}."
        )
        self.ext = ext


def extract_text(file_path: Path) -> tuple[str, bool, dict[str, Any]]:
    """Extract text + scanned-detect + metadata từ 1 file.

    Args:
        file_path: Path tuyệt đối hoặc tương đối đến file.

    Returns:
        text: str — full text extracted (\\n separator giữa pages/paragraphs).
        is_scanned: bool — True nếu PDF scanned (heuristic > 80% page < 30 ký tự).
        metadata: dict — {pages: int, format: str, encoding: str, ...}.

    Raises:
        UnsupportedFormatError: ext không trong ALLOWED_EXTENSIONS.
        FileNotFoundError: file không tồn tại.
    """
    if not file_path.exists():
        raise FileNotFoundError(f"File không tồn tại: {file_path}")
    ext = file_path.suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise UnsupportedFormatError(ext)

    if ext == ".docx":
        return _extract_docx(file_path)
    if ext in (".txt", ".md"):
        return _extract_text_file(file_path, ext)
    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".csv":
        return _extract_csv(file_path)
    if ext == ".xlsx":
        return _extract_xlsx(file_path)
    if ext == ".pptx":
        return _extract_pptx(file_path)
    if ext == ".html":
        return _extract_html(file_path)
    # Defensive — không reach do guard ALLOWED_EXTENSIONS phía trên.
    raise UnsupportedFormatError(ext)


def detect_scanned_pdf(file_path: Path) -> bool:
    """Public wrapper cho heuristic _is_pdf_scanned.

    Plan 04-04 BLOCKER #3 mitigation: router import detect_scanned_pdf để
    reject scanned PDF SYNCHRONOUSLY trước khi queue cocoindex flow
    (tránh persist document status='processing' rồi flow fail muộn).

    Args:
        file_path: Path tới PDF file.

    Returns:
        True nếu PDF scanned (> 80% page < 30 ký tự text).
        False nếu PDF text-only hoặc file không có page nào (defensive False).

    Raises:
        FileNotFoundError: file không tồn tại.
        UnsupportedFormatError: file không phải .pdf.
    """
    if not file_path.exists():
        raise FileNotFoundError(f"File không tồn tại: {file_path}")
    if file_path.suffix.lower() != ".pdf":
        raise UnsupportedFormatError(file_path.suffix.lower())
    reader = PdfReader(str(file_path))
    return _is_pdf_scanned(reader)


def _iter_block_items(parent: Any) -> Iterator[Paragraph | Table]:
    """Yield Paragraph / Table theo đúng thứ tự xuất hiện trong document body.

    python-docx KHÔNG expose thứ tự xen kẽ paragraph↔table sẵn — phải duyệt
    XML `<w:body>` children. Cần thiết để "tái hiện trung thực cấu trúc"
    (core value M2): tài liệu nội bộ Medinet thường xen kẽ đoạn văn + bảng.
    """
    body = parent.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, parent)
        elif child.tag == qn("w:tbl"):
            yield Table(child, parent)


def _table_to_text(table: Table) -> str:
    """Table → text: mỗi hàng = các ô join ' | ', các hàng join '\\n'.

    Đệ quy vào bảng lồng trong ô (`_Cell.text` KHÔNG bao gồm bảng lồng). Ô
    gộp (merged cell) lặp text — chấp nhận, đa số bảng layout tài liệu Medinet
    là 1×1 / không gộp.
    """
    lines: list[str] = []
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            cell_text = cell.text.strip()
            for nested in cell.tables:  # bảng lồng — gom thêm
                nested_text = _table_to_text(nested)
                if nested_text:
                    cell_text = f"{cell_text}\n{nested_text}".strip()
            if cell_text:
                cells.append(cell_text)
        line = " | ".join(cells)
        if line.strip():
            lines.append(line)
    return "\n".join(lines)


def _header_footer_parts(doc: Any) -> list[str]:
    """Text header/footer mọi section — dedupe nội bộ (header lặp mỗi section)."""
    seen: set[str] = set()
    parts: list[str] = []
    for section in doc.sections:
        for hf in (section.header, section.footer):
            for para in hf.paragraphs:
                t = para.text.strip()
                if t and t not in seen:
                    seen.add(t)
                    parts.append(t)
            for tbl in hf.tables:
                t = _table_to_text(tbl)
                if t and t not in seen:
                    seen.add(t)
                    parts.append(t)
    return parts


def _extract_docx(file_path: Path) -> tuple[str, bool, dict[str, Any]]:
    """Extract DOCX qua python-docx — duyệt body theo đúng thứ tự tài liệu.

    Đọc paragraph + table cell (gồm bảng lồng) theo thứ tự xuất hiện, cộng
    header/footer (dedupe vì lặp mỗi section). Tài liệu nội bộ Medinet thường
    dựng layout bằng bảng → bỏ bảng = mất nội dung (core value M2 yêu cầu tái
    hiện "bảng").
    """
    doc = DocxDocument(str(file_path))

    parts: list[str] = []
    paragraph_count = 0
    table_count = 0
    for block in _iter_block_items(doc):
        if isinstance(block, Table):
            table_count += 1
            block_text = _table_to_text(block)
            if block_text:
                parts.append(block_text)
        else:  # Paragraph
            para_text = block.text.strip()
            if para_text:
                paragraph_count += 1
                parts.append(para_text)

    parts.extend(_header_footer_parts(doc))
    text = "\n\n".join(parts)
    meta: dict[str, Any] = {
        "pages": 1,  # DOCX không có page count semantic — set 1 cho monotonic.
        "format": "docx",
        "encoding": "utf-8",
        "paragraph_count": paragraph_count,
        "table_count": table_count,
    }
    return text, False, meta


def _extract_text_file(file_path: Path, ext: str) -> tuple[str, bool, dict[str, Any]]:
    """Extract TXT/MD — UTF-8 trước, chardet fallback nếu fail."""
    raw = file_path.read_bytes()
    # Thử UTF-8 trước (case common cho VN medical doc trên web/docx-export).
    try:
        text = raw.decode("utf-8")
        encoding = "utf-8"
    except UnicodeDecodeError:
        # Fallback chardet detect (Windows-1258 / latin-1 / GB2312 / etc).
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "latin-1"
        text = raw.decode(encoding, errors="replace")
    meta: dict[str, Any] = {
        "pages": 1,
        "format": ext.lstrip("."),
        "encoding": encoding,
    }
    return text, False, meta


def _extract_pdf(file_path: Path) -> tuple[str, bool, dict[str, Any]]:
    """Extract PDF qua pypdf text-only + scanned-detect heuristic."""
    reader = PdfReader(str(file_path))
    pages_text: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        pages_text.append(page_text)
    text = "\n".join(pages_text)
    is_scanned = _is_pdf_scanned(reader)
    meta: dict[str, Any] = {
        "pages": len(reader.pages),
        "format": "pdf",
        "encoding": "utf-8",  # pypdf decode ra str unicode nội bộ.
    }
    return text, is_scanned, meta


def _is_pdf_scanned(reader: PdfReader) -> bool:
    """Heuristic: PDF scanned nếu > 80% page có < 30 ký tự text non-whitespace.

    pypdf trả empty/whitespace cho scanned PDF (R4 mitigation — gỡ Docling khỏi M2).
    Threshold 30 ký tự cover edge case "Trang 1" / page number tự render.

    Edge case: PDF 0 page → coi là scanned (defensive — Plan 04-04 sẽ trả 415
    với cùng error code `failed_unsupported`).
    """
    if len(reader.pages) == 0:
        return True
    empty_pages = 0
    for page in reader.pages:
        text = page.extract_text() or ""
        if len(text.strip()) < _SCANNED_PAGE_CHAR_THRESHOLD:
            empty_pages += 1
    return (empty_pages / len(reader.pages)) > _SCANNED_RATIO_THRESHOLD


def _decode_with_fallback(raw: bytes) -> tuple[str, str]:
    """Decode bytes ưu tiên utf-8-sig (strip BOM), fallback chardet detect.

    Dùng chung cho CSV + HTML — 2 format hay export từ Excel/Word Windows
    có BOM `\\ufeff` đầu file. `utf-8-sig` codec stdlib tự strip BOM,
    KHÔNG cần codecs.BOM_UTF8 detect thủ công.
    """
    try:
        return raw.decode("utf-8-sig"), "utf-8-sig"
    except UnicodeDecodeError:
        detected = chardet.detect(raw)
        enc = detected.get("encoding") or "latin-1"
        return raw.decode(enc, errors="replace"), enc


def _extract_csv(file_path: Path) -> tuple[str, bool, dict[str, Any]]:
    """Extract CSV — auto-detect delimiter + encoding, mỗi row join ' | '.

    Pitfall mitigation:
    - BOM (Excel Windows export): utf-8-sig codec tự strip.
    - Delimiter VN locale: csv.Sniffer detect trong {',', ';', '\\t', '|'}; fallback ','.
    - Encoding non-utf8: chardet fallback (latin-1 / windows-1258 / etc).

    Output format: mỗi row hàng = các ô join ' | ', các hàng join '\\n' —
    KHỚP pattern `_table_to_text` DOCX để chunker downstream xử lý nhất quán.
    """
    raw = file_path.read_bytes()
    text_decoded, encoding = _decode_with_fallback(raw)

    # Sniff delimiter từ sample 4KB đầu file
    sample = text_decoded[:4096]
    try:
        dialect = _csv.Sniffer().sniff(sample, delimiters=",;\t|")
        delimiter = dialect.delimiter
    except _csv.Error:
        delimiter = ","

    reader = _csv.reader(_io.StringIO(text_decoded), delimiter=delimiter)
    rows: list[str] = []
    row_count = 0
    for row in reader:
        cells = [cell.strip() for cell in row if cell.strip()]
        if cells:
            rows.append(" | ".join(cells))
            row_count += 1
    text = "\n".join(rows)
    meta: dict[str, Any] = {
        "pages": 1,
        "format": "csv",
        "encoding": encoding,
        "row_count": row_count,
        "delimiter": delimiter,
    }
    return text, False, meta


def _extract_xlsx(file_path: Path) -> tuple[str, bool, dict[str, Any]]:
    """Extract XLSX qua openpyxl read_only + data_only.

    - `read_only=True`: streaming parse, memory-friendly cho file lớn.
    - `data_only=True`: trả value của cell (kể cả công thức đã eval); cell công
      thức chưa eval (mới create chưa save Excel) → None → skip.
    - Sheet header `--- Sheet: <name> ---` trước mỗi sheet để chunker phân biệt.
    - Row format ` | ` separator tương đương CSV + DOCX table.
    """
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        total_rows = 0
        sheet_count = 0
        for sheet in wb.worksheets:
            sheet_count += 1
            sheet_lines: list[str] = [f"--- Sheet: {sheet.title} ---"]
            for row in sheet.iter_rows(values_only=True):
                cells = [
                    str(cell).strip()
                    for cell in row
                    if cell is not None and str(cell).strip()
                ]
                if cells:
                    sheet_lines.append(" | ".join(cells))
                    total_rows += 1
            if len(sheet_lines) > 1:  # có nội dung ngoài header
                parts.append("\n".join(sheet_lines))
    finally:
        wb.close()

    text = "\n\n".join(parts)
    meta: dict[str, Any] = {
        "pages": 1,
        "format": "xlsx",
        "encoding": "utf-8",
        "sheet_count": sheet_count,
        "row_count": total_rows,
    }
    return text, False, meta


def _extract_pptx(file_path: Path) -> tuple[str, bool, dict[str, Any]]:
    """Extract PPTX qua python-pptx — text frame mỗi slide + heading separator.

    - Duyệt `shape.has_text_frame` → `text_frame.text` (paragraph text gộp).
    - Slide header `--- Slide N ---` để chunker phân biệt biên slide.
    - Bỏ qua: Picture, Chart, Table shape (rare ở docs nội bộ; defer v4.0).
    """
    from pptx import Presentation

    prs = Presentation(str(file_path))
    parts: list[str] = []
    slide_count = 0
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_count += 1
        slide_lines: list[str] = [f"--- Slide {slide_idx} ---"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf_text = shape.text_frame.text.strip()
            if tf_text:
                slide_lines.append(tf_text)
        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))

    text = "\n\n".join(parts)
    meta: dict[str, Any] = {
        "pages": slide_count,
        "format": "pptx",
        "encoding": "utf-8",
        "slide_count": slide_count,
    }
    return text, False, meta


def _extract_html(file_path: Path) -> tuple[str, bool, dict[str, Any]]:
    """Extract HTML qua BeautifulSoup + lxml parser — strip tag, decompose script/style.

    Pitfall mitigation:
    - `<script>` + `<style>` tag content KHÔNG phải nội dung user → decompose trước get_text.
    - lxml parser handle HTML malformed tốt hơn html.parser stdlib (Word/PowerPoint export
      HTML hay sai cấu trúc).
    - BOM: utf-8-sig codec strip.
    """
    from bs4 import BeautifulSoup

    raw = file_path.read_bytes()
    text_decoded, encoding = _decode_with_fallback(raw)
    soup = BeautifulSoup(text_decoded, "lxml")

    # Decompose script + style trước khi get_text (T-quick-html-01 mitigation)
    for tag in soup(["script", "style"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    meta: dict[str, Any] = {
        "pages": 1,
        "format": "html",
        "encoding": encoding,
    }
    return text, False, meta
